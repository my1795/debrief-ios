#!/usr/bin/env python3
"""
Debrief Pitch Deck Generator
Creates a professional PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Colors
TEAL_900 = RGBColor(19, 78, 74)
TEAL_800 = RGBColor(17, 94, 89)
TEAL_500 = RGBColor(20, 184, 166)
TEAL_400 = RGBColor(45, 212, 191)
TEAL_300 = RGBColor(94, 234, 212)
EMERALD_900 = RGBColor(6, 78, 59)
WHITE = RGBColor(255, 255, 255)
GRAY_400 = RGBColor(156, 163, 175)
GRAY_700 = RGBColor(55, 65, 81)


def add_gradient_background(slide, prs):
    """Add a dark teal background to the slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = TEAL_900
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, prs)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = TEAL_300
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points, subtitle=None):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, prs)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle if provided
    y_offset = 1.3
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = TEAL_300
        y_offset = 1.8

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
        p.level = 0

    return slide


def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, prs)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    return slide


def add_metrics_slide(prs, title, metrics):
    """Add a slide with key metrics"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, prs)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Metrics boxes
    num_metrics = len(metrics)
    box_width = 2.8
    total_width = num_metrics * box_width + (num_metrics - 1) * 0.2
    start_x = (10 - total_width) / 2

    for i, (label, value) in enumerate(metrics):
        x = start_x + i * (box_width + 0.2)

        # Box background
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.5),
            Inches(box_width), Inches(2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = TEAL_800
        box.line.fill.background()

        # Value
        value_box = slide.shapes.add_textbox(Inches(x), Inches(2.7), Inches(box_width), Inches(1))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEAL_300
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(box_width), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide


def create_pitch_deck():
    """Create the complete pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Cover
    slide = add_title_slide(prs, "DEBRIEF", '"Never forget what matters from any call"')
    # Add tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Post-Call Personal Memory Assistant for iOS"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_400
    p.alignment = PP_ALIGN.CENTER

    # SLIDE 2: The Problem
    add_content_slide(prs, "The Problem", [
        "Professionals make 8-12 business calls daily",
        "",
        "50% of call details forgotten within 10 minutes",
        "",
        "90% forgotten within 24 hours",
        "",
        '"I\'ll take notes later" -> Never happens',
        "",
        "Result: Weak follow-ups, lost opportunities, damaged relationships"
    ], "Every Day, $37B Worth of Human Memory is Lost")

    # SLIDE 3: The Solution
    add_content_slide(prs, "The Solution", [
        "Call Ends -> Notification Appears -> Record Voice Memo",
        "",
        "AI generates: Summary + Action Items + Person-linked Memory",
        "",
        "HOW IT WORKS:",
        "  1. Phone call ends (0 sec)",
        "  2. Notification: 'Record debrief?' (2 sec)",
        "  3. User records voice memo (30 sec)",
        "  4. AI processes transcript + summary (60 sec)",
        "  5. Context available before next call (forever)",
        "",
        "PRIVACY-FIRST: No call recording. Only post-call voice memos."
    ], "Debrief: Automatic Memory Assistant When Calls End")

    # SLIDE 4: Product Demo
    add_two_column_slide(prs, "Product Demo",
        [
            "CORE FEATURES:",
            "",
            "CallKit Integration",
            "  -> Auto-trigger when calls end",
            "",
            "Voice-First Input",
            "  -> Zero typing, just speak",
            "",
            "AI Processing",
            "  -> Whisper + GPT-4",
            "",
            "Person-Centric",
            "  -> Memory linked to contacts",
        ],
        [
            "KEY SCREENS:",
            "",
            "Debriefs List",
            "  -> All past conversations",
            "",
            "Recording Flow",
            "  -> Auto-start, contact select",
            "",
            "Detail View",
            "  -> Summary, actions, audio",
            "",
            "Semantic Search",
            "  -> 'Budget discussions'"
        ]
    )

    # SLIDE 5: Market Opportunity
    add_metrics_slide(prs, "Market Opportunity", [
        ("TAM", "$21B"),
        ("SAM", "$3.8B"),
        ("SOM", "$380M")
    ])
    # Add additional context
    slide = prs.slides[-1]
    context_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    tf = context_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI Transcription: $4.5B -> $19.2B (2024-2034, CAGR 15.6%)"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_400
    p = tf.add_paragraph()
    p.text = "Speech-to-Text API: $5B -> $21B (2024-2034, CAGR 15.2%)"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_400
    p = tf.add_paragraph()
    p.text = "AI Meeting Assistant: $3.86B -> $29.45B (2025-2034, CAGR 25.62%)"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_400

    # SLIDE 6: Competitive Landscape
    add_two_column_slide(prs, "Competitive Landscape",
        [
            "COMPETITORS:",
            "",
            "Otter.ai",
            "  -> Web/Mobile, Manual trigger",
            "  -> Records actual calls",
            "",
            "Fireflies.ai",
            "  -> Web, Calendar bot",
            "  -> Enterprise focus",
            "",
            "Fathom",
            "  -> Web, Meeting bot",
            "  -> Zoom teams focus"
        ],
        [
            "DEBRIEF DIFFERENCE:",
            "",
            "Call-Triggered",
            "  -> Not manual, not bot-based",
            "",
            "Phone Calls Focus",
            "  -> Not meetings",
            "",
            "Privacy-First",
            "  -> NO call recording",
            "",
            "Mobile-Native",
            "  -> iOS optimized experience",
            "",
            "Zero Friction",
            "  -> 30 seconds, voice only"
        ]
    )

    # SLIDE 7: Business Model
    add_content_slide(prs, "Business Model", [
        "SUBSCRIPTION TIERS:",
        "",
        "FREE ($0/month)",
        "  -> 50 debriefs/week, 30 min/week, 500MB storage",
        "",
        "PERSONAL ($9.99/month)",
        "  -> Unlimited debriefs, 150 min/week, Unlimited storage",
        "",
        "PRO ($19.99/month)",
        "  -> Unlimited everything, Priority AI, Advanced search",
        "",
        "TARGET UNIT ECONOMICS:",
        "  -> CAC: $15-25 | LTV: $150+ | LTV:CAC: 6:1+"
    ])

    # SLIDE 8: PLG Strategy
    add_content_slide(prs, "Product-Led Growth Strategy", [
        "PLG FLYWHEEL:",
        "",
        "1. User Activation",
        "   -> Call ends, record debrief in <30 seconds",
        "",
        "2. Habit Formation",
        "   -> Daily debrief routine, call-triggered engagement",
        "",
        "3. Value Realization",
        "   -> Better follow-ups, stronger relationships",
        "",
        "4. Organic Expansion",
        "   -> Share insights, hit quota, upgrade tier",
        "",
        "AHA MOMENT: See first AI summary (<2 minutes from call end)"
    ])

    # SLIDE 9: Technology
    add_two_column_slide(prs, "Technology Stack",
        [
            "iOS CLIENT:",
            "",
            "SwiftUI + MVVM",
            "  -> Modern, maintainable",
            "",
            "CallKit Integration",
            "  -> Native call detection",
            "",
            "AVFoundation",
            "  -> High-quality recording",
            "",
            "Offline-First",
            "  -> Record anywhere, sync later"
        ],
        [
            "BACKEND + AI:",
            "",
            "Firebase / GCP",
            "  -> Firestore, Storage, Cloud Run",
            "",
            "OpenAI Whisper",
            "  -> Transcription",
            "",
            "GPT-4",
            "  -> Summary + Action Items",
            "",
            "Vector Embeddings",
            "  -> Semantic Search",
            "",
            "AES-256 E2E Encryption"
        ]
    )

    # SLIDE 10: Traction & Milestones
    add_two_column_slide(prs, "Traction & Milestones",
        [
            "COMPLETED:",
            "",
            "[x] iOS MVP Complete",
            "",
            "[x] CallKit Integration",
            "",
            "[x] AI Pipeline (Whisper + GPT-4)",
            "",
            "[x] Billing System v2",
            "",
            "[x] Semantic Search",
            "",
            "[x] TestFlight Beta Live",
            "",
            "[x] App Store Ready"
        ],
        [
            "ROADMAP:",
            "",
            "Q1 2026",
            "  -> iOS App Store Launch",
            "",
            "Q2 2026",
            "  -> 10K MAU",
            "",
            "Q3 2026",
            "  -> Android Beta",
            "",
            "Q4 2026",
            "  -> 100K MAU"
        ]
    )

    # SLIDE 11: The Ask
    add_metrics_slide(prs, "The Ask: $1M Seed Round", [
        ("Engineering", "50%"),
        ("Growth", "30%"),
        ("Operations", "20%")
    ])
    slide = prs.slides[-1]
    details_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    tf = details_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Engineering: iOS Dev, Backend, ML/AI Engineer"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_400
    p = tf.add_paragraph()
    p.text = "Growth: User Acquisition, Content Marketing, Community"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_400
    p = tf.add_paragraph()
    p.text = "Operations: Infrastructure, Legal, Runway"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_400
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "MILESTONES: App Store M1 | 10K MAU M6 | Android M9 | $1M ARR M18"
    p.font.size = Pt(16)
    p.font.color.rgb = TEAL_300

    # SLIDE 12: Why Now
    add_content_slide(prs, "Why Now?", [
        "1. AI TRANSCRIPTION COST DOWN",
        "   Whisper API: $0.006/min (2x cheaper than 2023)",
        "",
        "2. MOBILE AI CAPABILITIES UP",
        "   On-device ML, faster processing, better UX",
        "",
        "3. REMOTE WORK = MORE CALLS",
        "   Phone calls up 40% since 2020",
        "",
        "4. MEETING FATIGUE -> CALL PREFERENCE",
        "   Users prefer quick calls over scheduled meetings",
        "",
        "5. PRIVACY AWARENESS UP",
        "   Users want control, not surveillance"
    ], "Perfect Market Timing")

    # SLIDE 13: Closing
    slide = add_title_slide(prs, "DEBRIEF", "Remember every call. Follow up better.\nBuild stronger relationships.")

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
    tf = points_box.text_frame
    tf.word_wrap = True
    points = [
        "Call-Triggered (not meeting-bot)",
        "Privacy-First (no call recording)",
        "Mobile-Native (iOS optimized)",
        "AI-Powered (Whisper + GPT-4)",
        "PLG (self-service growth)"
    ]
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEAL_300
        p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "Debrief_Pitch_Deck.pptx")
    prs.save(output_path)
    print(f"\n{'='*60}")
    print(f"  PITCH DECK CREATED SUCCESSFULLY!")
    print(f"{'='*60}")
    print(f"\n  File: {output_path}")
    print(f"\n  Slides: 13")
    print(f"  1. Cover")
    print(f"  2. The Problem")
    print(f"  3. The Solution")
    print(f"  4. Product Demo")
    print(f"  5. Market Opportunity")
    print(f"  6. Competitive Landscape")
    print(f"  7. Business Model")
    print(f"  8. PLG Strategy")
    print(f"  9. Technology Stack")
    print(f" 10. Traction & Milestones")
    print(f" 11. The Ask")
    print(f" 12. Why Now")
    print(f" 13. Closing")
    print(f"\n{'='*60}")
    print(f"  Open with Keynote or PowerPoint")
    print(f"{'='*60}\n")
    return output_path


if __name__ == "__main__":
    create_pitch_deck()
